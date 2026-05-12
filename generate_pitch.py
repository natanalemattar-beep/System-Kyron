from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BG_DARK    = RGBColor(0x04, 0x06, 0x0F)
ACCENT     = RGBColor(0x3B, 0x82, 0xF6)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
WHITE_DIM  = RGBColor(0x94, 0xA3, 0xB8)
GREEN      = RGBColor(0x10, 0xB9, 0x81)
AMBER      = RGBColor(0xF5, 0x9E, 0x0B)
ROSE       = RGBColor(0xF4, 0x3F, 0x5E)

blank_layout = prs.slide_layouts[6]

def add_slide():
    return prs.slides.add_slide(blank_layout)

def fill_bg(slide, color=BG_DARK):
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, fill_color, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    return shape

def add_text(slide, text, l, t, w, h, size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Segoe UI"
    return txBox

def add_bullet_box(slide, bullets, l, t, w, h, size=14, color=WHITE_DIM):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for bullet in bullets:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = bullet
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = "Segoe UI"
    return txBox

# ─────────────────────────────────────────────
# SLIDE 1 — PORTADA
# ─────────────────────────────────────────────
s = add_slide()
fill_bg(s)
add_rect(s, 0, 0, 13.33, 0.08, ACCENT)           # top bar
add_rect(s, 0, 7.42, 13.33, 0.08, ACCENT)         # bottom bar
add_text(s, "RETO INSPIRA 2026", 0.6, 0.3, 12, 0.5, size=11, color=ACCENT, bold=True)
add_text(s, "SYSTEM KYRON", 0.6, 1.2, 10, 2.2, size=72, bold=True, color=WHITE)
add_text(s, "Todo lo que necesitas para vender y crecer:", 0.6, 3.6, 10, 0.6, size=22, color=WHITE_DIM)
add_text(s, "tus líneas, tu web y cero complicaciones.", 0.6, 4.1, 10, 0.6, size=22, color=WHITE_DIM, italic=True)
add_rect(s, 0.6, 5.0, 3.5, 0.05, ACCENT)
add_text(s, "Carlos Mattar — Fundador & CEO", 0.6, 5.2, 8, 0.45, size=14, color=WHITE_DIM)
add_text(s, "Emprendimiento Carlos Mattar  |  RIF J-50832149-9", 0.6, 5.7, 8, 0.45, size=12, color=RGBColor(0x47,0x55,0x69))
add_text(s, "📱 0424-1846016  |  @systemkyron  |  system-kyron.vercel.app", 0.6, 6.2, 10, 0.45, size=11, color=RGBColor(0x47,0x55,0x69))

# ─────────────────────────────────────────────
# SLIDE 2 — EL PROBLEMA
# ─────────────────────────────────────────────
s = add_slide()
fill_bg(s)
add_rect(s, 0, 0, 13.33, 0.08, ROSE)
add_rect(s, 0, 7.42, 13.33, 0.08, ROSE)
add_text(s, "02  |  DEFINICIÓN DEL PROBLEMA", 0.6, 0.25, 12, 0.4, size=11, color=ROSE, bold=True)
add_text(s, "El Dolor de Cabeza Real", 0.6, 0.8, 8, 1.0, size=38, bold=True, color=WHITE)
add_text(s, "Pelear con la tecnología no debería ser tu trabajo.", 0.6, 1.9, 9, 0.6, size=18, color=WHITE_DIM, italic=True)

body = (
    "Los emprendedores y pequeños empresarios pierden tiempo y dinero contratando entre "
    "3 y 5 proveedores distintos para resolver cosas básicas: instalar líneas, armar su web "
    "y organizar sus procesos.\n\n"
    "Más del 60% de los emprendedores venezolanos no tiene presencia digital activa. "
    "Quienes la tienen reportan que gestionar varios proveedores les consume hasta un 40% "
    "de su tiempo productivo semanal.\n\n"
    "El resultado: negocios que no venden todo lo que podrían porque están atrapados "
    "resolviendo tecnología en vez de atender a sus clientes."
)
add_text(s, body, 0.6, 2.6, 7.8, 4.0, size=14, color=WHITE_DIM)

# Stats box
for i, (val, lbl) in enumerate([("60%","Sin presencia digital"),("40%","Tiempo perdido/semana"),("+4","Proveedores separados")]):
    x = 9.0 + i * 1.4 if i < 2 else 9.0
    y = 2.5 + (0 if i < 2 else 2.1)
    w = 3.8 if i == 2 else 1.9
    add_rect(s, 9.0 if i==0 else 11.1 if i==1 else 9.0,
             2.5 if i<2 else 4.8, 
             1.9 if i<2 else 3.8,
             1.7, RGBColor(0x1E,0x1E,0x2E))
    xb = 9.1 if i==0 else 11.2 if i==1 else 9.1
    yb = 2.6 if i<2 else 4.9
    add_text(s, val, xb, yb, 1.7 if i<2 else 3.6, 0.8, size=32, bold=True, color=ROSE, align=PP_ALIGN.CENTER)
    add_text(s, lbl, xb, yb+0.8, 1.7 if i<2 else 3.6, 0.6, size=10, color=WHITE_DIM, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────
# SLIDE 3 — PROPUESTA DE VALOR
# ─────────────────────────────────────────────
s = add_slide()
fill_bg(s)
add_rect(s, 0, 0, 13.33, 0.08, ACCENT)
add_rect(s, 0, 7.42, 13.33, 0.08, ACCENT)
add_text(s, "03  |  PROPUESTA DE VALOR", 0.6, 0.25, 12, 0.4, size=11, color=ACCENT, bold=True)
add_text(s, "La Solución: Todo en Un Solo Lugar", 0.6, 0.8, 9, 1.0, size=34, bold=True, color=WHITE)
add_text(s, "System Kyron es la única empresa que te da todo lo que necesita tu negocio, sin complicaciones.", 0.6, 1.85, 9, 0.6, size=16, color=WHITE_DIM, italic=True)

services = [
    ("📞  Líneas Corporativas", "Comunicación para todo tu equipo, activa al instante y sin contratos imposibles de entender."),
    ("🌐  Página Web Premium", "Tu plataforma web hecha a medida, rápida, hermosa y disponible las 24 horas del día."),
    ("⚙️  Soluciones Sostenibles", "Herramientas digitales para organizar tu negocio, eliminar el papel y crecer sin caos."),
]
for i, (title, desc) in enumerate(services):
    y = 2.7 + i * 1.4
    add_rect(s, 0.6, y, 7.8, 1.2, RGBColor(0x0D,0x11,0x1E))
    add_text(s, title, 0.9, y+0.1, 7.2, 0.4, size=15, bold=True, color=ACCENT)
    add_text(s, desc,  0.9, y+0.55, 7.2, 0.55, size=13, color=WHITE_DIM)

add_text(s, "¿Qué nos hace únicos?", 9.0, 2.5, 3.8, 0.4, size=13, bold=True, color=WHITE)
diffs = ["✅ Un solo proveedor para todo", "✅ Precios claros y transparentes", "✅ Planes que crecen contigo", "✅ Atención humana, no robots", "✅ Formalización legal incluida"]
add_bullet_box(s, diffs, 9.0, 3.0, 3.9, 3.5, size=13, color=WHITE_DIM)

# ─────────────────────────────────────────────
# SLIDE 4 — MERCADO OBJETIVO
# ─────────────────────────────────────────────
s = add_slide()
fill_bg(s)
add_rect(s, 0, 0, 13.33, 0.08, GREEN)
add_rect(s, 0, 7.42, 13.33, 0.08, GREEN)
add_text(s, "04  |  MERCADO OBJETIVO", 0.6, 0.25, 12, 0.4, size=11, color=GREEN, bold=True)
add_text(s, "¿A Quién Le Hablamos?", 0.6, 0.8, 9, 0.9, size=36, bold=True, color=WHITE)

profile = [
    "👤  Emprendedores, dueños de negocios familiares y directores de PyMEs",
    "📍  Venezuela (con proyección regional a corto plazo)",
    "🎂  Edad: 22 a 50 años",
    "💡  Intereses: Crecer su negocio, ahorrar tiempo, imagen profesional",
    "📱  Comportamiento: Activos en redes, buscan soluciones rápidas y de confianza",
]
add_text(s, "Perfil del Cliente Ideal", 0.6, 1.85, 7, 0.4, size=15, bold=True, color=GREEN)
add_bullet_box(s, profile, 0.6, 2.4, 7.5, 3.0, size=14, color=WHITE_DIM)

add_text(s, "Tamaño del Mercado", 8.8, 1.85, 4.1, 0.4, size=15, bold=True, color=GREEN)
market = [
    ("500K+", "Micro y PyMEs activas en Venezuela"),
    ("5,000", "Clientes potenciales con solo el 1% del mercado"),
    ("Regional", "Potencial de expansión a países similares"),
]
for i, (val, lbl) in enumerate(market):
    y = 2.5 + i * 1.6
    add_rect(s, 8.8, y, 4.1, 1.3, RGBColor(0x0D,0x11,0x1E))
    add_text(s, val, 8.9, y+0.05, 3.9, 0.65, size=24, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    add_text(s, lbl, 8.9, y+0.7, 3.9, 0.5, size=11, color=WHITE_DIM, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────
# SLIDE 5 — MODELO DE NEGOCIO
# ─────────────────────────────────────────────
s = add_slide()
fill_bg(s)
add_rect(s, 0, 0, 13.33, 0.08, AMBER)
add_rect(s, 0, 7.42, 13.33, 0.08, AMBER)
add_text(s, "05  |  MODELO DE NEGOCIO", 0.6, 0.25, 12, 0.4, size=11, color=AMBER, bold=True)
add_text(s, "¿Cómo Generamos Ingresos?", 0.6, 0.8, 9, 0.9, size=36, bold=True, color=WHITE)
add_text(s, "Simple, transparente y escalable.", 0.6, 1.75, 7, 0.5, size=17, color=WHITE_DIM, italic=True)

streams = [
    ("01", "Suscripción — Líneas Corporativas", "Cobro mensual recurrente por cada línea de comunicación activa en la empresa del cliente."),
    ("02", "Proyecto + Mantenimiento — Webs",   "Pago inicial por el desarrollo web más una cuota mensual de mantenimiento y actualizaciones."),
    ("03", "Paquetes Integrales",                "Planes combinados (líneas + web + soluciones) con descuento. El cliente ahorra, nosotros fidelizamos."),
]
for i, (num, title, desc) in enumerate(streams):
    y = 2.5 + i * 1.55
    add_rect(s, 0.6, y, 11.8, 1.35, RGBColor(0x0D,0x11,0x1E))
    add_text(s, num, 0.75, y+0.1, 0.7, 1.1, size=28, bold=True, color=AMBER, align=PP_ALIGN.CENTER)
    add_rect(s, 1.55, y+0.2, 0.04, 0.9, AMBER)
    add_text(s, title, 1.75, y+0.1, 8.0, 0.5, size=15, bold=True, color=WHITE)
    add_text(s, desc,  1.75, y+0.62, 8.0, 0.6, size=13, color=WHITE_DIM)

add_text(s, "Escalabilidad: A mayor número de clientes, menor costo por cliente → mayor margen de ganancia.", 0.6, 7.1, 11.8, 0.35, size=11, color=RGBColor(0x47,0x55,0x69), italic=True)

# ─────────────────────────────────────────────
# SLIDE 6 — MARKETING Y VENTAS
# ─────────────────────────────────────────────
s = add_slide()
fill_bg(s)
add_rect(s, 0, 0, 13.33, 0.08, RGBColor(0x8B,0x5C,0xF6))
add_rect(s, 0, 7.42, 13.33, 0.08, RGBColor(0x8B,0x5C,0xF6))
add_text(s, "06  |  ESTRATEGIA DE MARKETING Y VENTAS", 0.6, 0.25, 12, 0.4, size=11, color=RGBColor(0x8B,0x5C,0xF6), bold=True)
add_text(s, "¿Cómo Nos Van a Conocer?", 0.6, 0.8, 9, 0.9, size=36, bold=True, color=WHITE)

channels = [
    ("📱  Redes Sociales (Instagram & TikTok)",
     "Contenido educativo en el lenguaje del emprendedor: 'Cómo armar tu web sin pagar de más', 'Por qué necesitas líneas corporativas'. Generamos confianza antes de vender."),
    ("🤝  Referidos y Boca a Boca",
     "Cada cliente satisfecho nos recomienda. Programa de referidos: el cliente gana un descuento por cada empresa que nos trae."),
    ("🏛️  Alianzas con Incubadoras y Cámaras de Comercio",
     "Nos posicionamos como el proveedor tecnológico recomendado para nuevos negocios en proceso de formalización."),
    ("🌐  Demostración en Vivo — nuestra propia web",
     "system-kyron.vercel.app es nuestra mejor carta de presentación. El cliente ve el producto real antes de comprar."),
]
for i, (ch, desc) in enumerate(channels):
    col = i % 2
    row = i // 2
    x = 0.6 + col * 6.4
    y = 2.1 + row * 2.4
    add_rect(s, x, y, 6.1, 2.1, RGBColor(0x0D,0x11,0x1E))
    add_text(s, ch, x+0.2, y+0.12, 5.7, 0.5, size=13, bold=True, color=RGBColor(0x8B,0x5C,0xF6))
    add_text(s, desc, x+0.2, y+0.65, 5.7, 1.3, size=12, color=WHITE_DIM)

# ─────────────────────────────────────────────
# SLIDE 7 — IMPACTO SOCIAL Y AMBIENTAL
# ─────────────────────────────────────────────
s = add_slide()
fill_bg(s)
add_rect(s, 0, 0, 13.33, 0.08, GREEN)
add_rect(s, 0, 7.42, 13.33, 0.08, GREEN)
add_text(s, "07  |  IMPACTO SOCIAL Y AMBIENTAL", 0.6, 0.25, 12, 0.4, size=11, color=GREEN, bold=True)
add_text(s, "Más que un Negocio, un Cambio Real", 0.6, 0.8, 10, 0.9, size=34, bold=True, color=WHITE)

add_text(s, "🤝  Impacto Social", 0.6, 2.0, 5.8, 0.5, size=18, bold=True, color=GREEN)
social_text = (
    "System Kyron democratiza el acceso a la tecnología. "
    "Un emprendedor de barrio tiene hoy la misma capacidad de proyección digital "
    "que una empresa grande, sin necesitar grandes presupuestos.\n\n"
    "Esto genera empleos formales, fortalece la economía local y ayuda a más venezolanos "
    "a formalizarse y a crecer."
)
add_text(s, social_text, 0.6, 2.6, 5.8, 3.5, size=14, color=WHITE_DIM)

add_rect(s, 6.7, 1.8, 0.05, 5.2, RGBColor(0x1E,0x29,0x3B))

add_text(s, "🌱  Impacto Ambiental", 7.0, 2.0, 5.8, 0.5, size=18, bold=True, color=GREEN)
env_text = (
    "Promovemos el modelo 100% cero papel. "
    "Todos nuestros procesos son digitales: contratos, facturas, comunicaciones y gestión.\n\n"
    "Cada cliente que sumamos al sistema digital deja de usar papel en sus operaciones diarias, "
    "reduciendo su huella de carbono desde el primer día que nos contrata."
)
add_text(s, env_text, 7.0, 2.6, 5.8, 3.5, size=14, color=WHITE_DIM)

add_text(s, "\"No vendemos solo tecnología. Vendemos un futuro más organizado, más justo y más sostenible.\"", 0.6, 6.7, 11.8, 0.5, size=12, color=RGBColor(0x47,0x55,0x69), italic=True, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────
# SLIDE 8 — ESTADO ACTUAL Y HOJA DE RUTA
# ─────────────────────────────────────────────
s = add_slide()
fill_bg(s)
add_rect(s, 0, 0, 13.33, 0.08, ACCENT)
add_rect(s, 0, 7.42, 13.33, 0.08, ACCENT)
add_text(s, "08  |  ESTADO ACTUAL Y HOJA DE RUTA", 0.6, 0.25, 12, 0.4, size=11, color=ACCENT, bold=True)
add_text(s, "¿Dónde Estamos y a Dónde Vamos?", 0.6, 0.8, 10, 0.9, size=33, bold=True, color=WHITE)

add_text(s, "Estado actual: Prototipo Funcional Avanzado con clientes reales", 0.6, 1.8, 11, 0.5, size=15, color=GREEN, bold=True)
status = [
    "✅  Plataforma web en producción y funcionando",
    "✅  Negocio registrado legalmente (RIF J-50832149-9)",
    "✅  Primeros clientes activos en líneas y webs",
]
add_bullet_box(s, status, 0.6, 2.35, 11.5, 1.2, size=13, color=WHITE_DIM)

add_text(s, "Objetivos — Próximos 6 Meses", 0.6, 3.7, 11, 0.45, size=15, bold=True, color=ACCENT)

milestones = [
    ("MES 1-2", "LANZAMIENTO", "Campaña activa en redes. Alcanzar los primeros 50 clientes pagos entre líneas y webs.", ROSE),
    ("MES 3-4", "CRECIMIENTO",  "3 alianzas formales con cámaras o incubadoras. Alcanzar 100 clientes activos.", AMBER),
    ("MES 5-6", "CONSOLIDACIÓN","Lanzar paquetes integrales. Lograr ingresos mensuales recurrentes autosustentables.", GREEN),
]
for i, (period, phase, desc, color) in enumerate(milestones):
    x = 0.6 + i * 4.2
    y = 4.3
    add_rect(s, x, y, 3.9, 2.7, RGBColor(0x0D,0x11,0x1E))
    add_text(s, period, x+0.15, y+0.15, 3.6, 0.35, size=10, bold=True, color=color)
    add_text(s, phase,  x+0.15, y+0.55, 3.6, 0.5, size=16, bold=True, color=WHITE)
    add_text(s, desc,   x+0.15, y+1.15, 3.6, 1.4, size=12, color=WHITE_DIM)
    add_rect(s, x, y+2.55, 3.9, 0.08, color)

# ─────────────────────────────────────────────
# SLIDE 9 — CIERRE / CONTACTO
# ─────────────────────────────────────────────
s = add_slide()
fill_bg(s)
add_rect(s, 0, 0, 13.33, 0.08, GREEN)
add_rect(s, 0, 7.42, 13.33, 0.08, GREEN)
add_text(s, "HAGÁMOSLO REALIDAD", 0.6, 1.2, 11, 1.5, size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "Deja que nosotros nos encarguemos de la tecnología.", 0.6, 2.9, 12.1, 0.6, size=20, color=WHITE_DIM, align=PP_ALIGN.CENTER, italic=True)
add_text(s, "Tú dedícate a hacer lo que amas y a llevar tu negocio al siguiente nivel.", 0.6, 3.5, 12.1, 0.6, size=18, color=WHITE_DIM, align=PP_ALIGN.CENTER, italic=True)
add_rect(s, 3.5, 4.5, 6.33, 0.05, GREEN)
add_text(s, "NUESTRA PROMESA:", 3.5, 4.7, 6.5, 0.4, size=12, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
promises = ["✔  Te hablamos claro, sin términos raros o técnicos.",
            "✔  Soporte rápido de personas reales, no robots.",
            "✔  Nos encargamos de todo el trabajo pesado por ti."]
add_bullet_box(s, promises, 3.0, 5.2, 7.3, 1.3, size=13, color=WHITE_DIM)
add_text(s, "📱 0424-1846016   |   📸 @systemkyron   |   🌐 system-kyron.vercel.app", 0.6, 6.8, 12.1, 0.45, size=12, color=RGBColor(0x47,0x55,0x69), align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────
# GUARDAR
# ─────────────────────────────────────────────
output_path = r"c:\Users\Carlos\.gemini\antigravity\scratch\System-Kyron-Git\Resumen_Ejecutivo_SystemKyron_RetoInspira.pptx"
prs.save(output_path)
print(f"✅ PPTX generado en: {output_path}")
